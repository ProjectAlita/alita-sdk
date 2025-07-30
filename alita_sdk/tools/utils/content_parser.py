import re

from docx import Document
from io import BytesIO
import pandas as pd
from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import io
import pymupdf
from langchain_core.tools import ToolException
from transformers import BlipProcessor, BlipForConditionalGeneration
from langchain_core.messages import HumanMessage
from logging import getLogger

from ...runtime.langchain.tools.utils import bytes_to_base64

logger = getLogger(__name__)

image_processing_prompt='''
You are an AI model designed for analyzing images. Your task is to accurately describe the content of the given image. Depending on the type of image, follow these specific instructions:

If the image is a diagram (e.g., chart, table, pie chart, bar graph, etc.):

Identify the type of diagram.
Extract all numerical values, labels, axis titles, headings, legends, and any other textual elements.
Describe the relationships or trends between the data, if visible.
If the image is a screenshot:

Describe what is shown in the screenshot.
If it is a software interface, identify the program or website name (if visible).
List the key interface elements (e.g., buttons, menus, text fields, images, headers).
If there is text, extract it.
If the screenshot shows a conversation, describe the participants, the content of the messages, and timestamps (if visible).
If the image is a photograph:

Describe the main objects, people, animals, or elements visible in the photo.
Specify the setting (e.g., indoors, outdoors, nature, urban area).
If possible, identify the actions being performed by people or objects in the photo.
If the image is an illustration or drawing:

Describe the style of the illustration (e.g., realistic, cartoonish, abstract).
Identify the main elements, their colors, and the composition of the image.
If there is text, extract it.
If the image contains text:

Extract all text from the image.
Specify the format of the text (e.g., heading, paragraph, list).
If the image is a mixed type (e.g., a diagram within a screenshot):

Identify all types of content present in the image.
Perform an analysis for each type of content separately, following the relevant instructions above.
If the image does not fit into any of the above categories:

Provide a detailed description of what is shown in the image.
Highlight any visible details that could help in understanding the image.
Be as precise and thorough as possible in your responses. If something is unclear or illegible, state that explicitly.
'''

IMAGE_EXTENSIONS = ['jpg', 'jpeg', 'png', 'gif', 'bmp', 'tiff', 'webp', 'svg']


def parse_file_content(file_name=None, file_content=None, is_capture_image: bool = False, page_number: int = None,
                       sheet_name: str = None, llm=None, file_path: str = None, excel_by_sheets: bool = False):
    """Parse the content of a file based on its type and return the parsed content.

    Args:
        file_name (str): The name of the file to parse.
        file_content (bytes): The content of the file as bytes.
        is_capture_image (bool): Whether to capture images from the file.
        page_number (int, optional): The specific page number to parse for PDF or PPTX files.
        sheet_name (str, optional): The specific sheet name to parse for Excel files.
        llm: The language model to use for image processing.
        file_path (str, optional): The path to the file if it needs to be read from disk.
    Returns:
        str: The parsed content of the file.
    Raises:
        ToolException: If the file type is not supported or if there is an error reading the file.
        """

    if (file_path and (file_name or file_content)) or (not file_path and (not file_name or file_content is None)):
        raise ToolException("Either (file_name and file_content) or file_path must be provided, but not both.")

    if file_path:
        file_content = file_to_bytes(file_path)
        if file_content is None:
            return ToolException(f"File not found or could not be read: {file_path}")
        file_name = file_path.split('/')[-1]  # Extract file name from path
    if file_name.endswith('.txt'):
        return parse_txt(file_content)
    elif file_name.endswith('.docx'):
        return read_docx_from_bytes(file_content)
    elif file_name.endswith('.xlsx') or file_name.endswith('.xls'):
        return parse_excel(file_content, sheet_name, excel_by_sheets)
    elif file_name.endswith('.pdf'):
        return parse_pdf(file_content, page_number, is_capture_image, llm)
    elif file_name.endswith('.pptx'):
        return parse_pptx(file_content, page_number, is_capture_image, llm)
    elif any(file_name.lower().endswith(f".{ext}") for ext in IMAGE_EXTENSIONS):
        match = re.search(r'\.([a-zA-Z0-9]+)$', file_name)
        return __perform_llm_prediction_for_image(llm, file_content, match.group(1), image_processing_prompt)
    else:
        return ToolException(
            "Not supported type of files entered. Supported types are TXT, DOCX, PDF, PPTX, XLSX and XLS only.")

def parse_txt(file_content):
    try:
        return file_content.decode('utf-8')
    except Exception as e:
        return ToolException(f"Error decoding file content: {e}")

def parse_excel(file_content, sheet_name = None, return_by_sheets: bool = False):
    try:
        excel_file = io.BytesIO(file_content)
        if sheet_name:
            return parse_sheet(excel_file, sheet_name)
        dfs = pd.read_excel(excel_file, sheet_name=sheet_name)

        if return_by_sheets:
            result = {}
            for sheet_name, df in dfs.items():
                df.fillna('', inplace=True)
                result[sheet_name] = df.to_dict(orient='records')
            return result
        else:
            result = []
            for sheet_name, df in dfs.items():
                df.fillna('', inplace=True)
                string_content = df.to_string(index=False)
                result.append(f"====== Sheet name: {sheet_name} ======\n{string_content}")
            return "\n\n".join(result)
    except Exception as e:
        return ToolException(f"Error reading Excel file: {e}")

def parse_sheet(excel_file, sheet_name):
    df = pd.read_excel(excel_file, sheet_name=sheet_name)
    df.fillna('', inplace=True)
    return df.to_string()

def parse_pdf(file_content, page_number, is_capture_image, llm):
    with pymupdf.open(stream=file_content, filetype="pdf") as report:
        text_content = ''
        if page_number is not None:
            page = report.load_page(page_number - 1)
            text_content += read_pdf_page(report, page, page_number, is_capture_image, llm)
        else:
            for index, page in enumerate(report, start=1):
                text_content += read_pdf_page(report, page, index, is_capture_image, llm)
        return text_content

def parse_pptx(file_content, page_number, is_capture_image, llm=None):
    prs = Presentation(io.BytesIO(file_content))
    text_content = ''
    if page_number is not None:
        text_content += read_pptx_slide(prs.slides[page_number - 1], page_number, is_capture_image, llm)
    else:
        for index, slide in enumerate(prs.slides, start=1):
            text_content += read_pptx_slide(slide, index, is_capture_image, llm)
    return text_content

def read_pdf_page(report, page, index, is_capture_images, llm=None):
    text_content = f'Page: {index}\n'
    text_content += page.get_text()
    if is_capture_images:
        images = page.get_images(full=True)
        for i, img in enumerate(images):
            xref = img[0]
            base_image = report.extract_image(xref)
            img_bytes = base_image["image"]
            text_content += __perform_llm_prediction_for_image(llm, img_bytes)
    return text_content

def read_docx_from_bytes(file_content):
    """Read and return content from a .docx file using a byte stream."""
    try:
        doc = Document(BytesIO(file_content))
        text = []
        for paragraph in doc.paragraphs:
            text.append(paragraph.text)
        return '\n'.join(text)
    except Exception as e:
        print(f"Error reading .docx from bytes: {e}")
        return ""

def read_pptx_slide(slide, index, is_capture_image, llm):
    text_content = f'Slide: {index}\n'
    for shape in slide.shapes:
        if hasattr(shape, "text"):
            text_content += shape.text + "\n"
        elif is_capture_image and shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            try:
                caption = __perform_llm_prediction_for_image(llm, shape.image.blob)
            except:
                caption = "\n[Picture: unknown]\n"
            text_content += caption
    return text_content

def describe_image(image):
    processor = BlipProcessor.from_pretrained("Salesforce/blip-image-captioning-base")
    model = BlipForConditionalGeneration.from_pretrained("Salesforce/blip-image-captioning-base")
    inputs = processor(image, return_tensors="pt")
    out = model.generate(**inputs)
    return "\n[Picture: " + processor.decode(out[0], skip_special_tokens=True) + "]\n"

def __perform_llm_prediction_for_image(llm, image: bytes, image_format='png', prompt=image_processing_prompt) -> str:
    if not llm:
        raise ToolException("LLM is not provided for image processing.")
    base64_string = bytes_to_base64(image)
    result = llm.invoke([
        HumanMessage(
            content=[
                {"type": "text", "text": prompt},
                {
                    "type": "image_url",
                    "image_url": {"url": f"data:image/{image_format};base64,{base64_string}"},
                },
            ])
    ])
    return f"\n[Image description: {result.content}]\n"

# TODO: review usage of this function alongside with functions above
def load_content(file_path: str, extension: str = None, loader_extra_config: dict = None, llm = None) -> str:
    """
    Loads the content of a file based on its extension using a configured loader.
    """
    try:
        from ...runtime.langchain.document_loaders.constants import loaders_map

        if not extension:
            extension = file_path.split('.')[-1].lower()

        loader_config = loaders_map.get(extension)
        if not loader_config:
            logger.warning(f"No loader found for file extension: {extension}. File: {file_path}")
            return ""

        loader_cls = loader_config['class']
        loader_kwargs = loader_config['kwargs']

        if loader_extra_config:
            loader_kwargs.update(loader_extra_config)
        if loader_config['is_multimodal_processing'] and llm:
            loader_kwargs.update({'llm': llm})

        loader = loader_cls(file_path, **loader_kwargs)
        documents = loader.load()

        page_contents = [doc.page_content for doc in documents]
        return "\n".join(page_contents)
    except Exception as e:
        error_message = f"Error loading attachment: {str(e)}"
        logger.warning(f"{error_message} for file {file_path}")
        return ""

def load_content_from_bytes(file_content: bytes, extension: str = None, loader_extra_config: dict = None, llm = None) -> str:
    """Loads the content of a file from bytes based on its extension using a configured loader."""

    import tempfile

    # Automatic cleanup with context manager
    with tempfile.NamedTemporaryFile(mode='w+b', delete=True) as temp_file:
        # Write data to temp file
        temp_file.write(file_content)
        temp_file.flush()  # Ensure data is written

        # Get the file path for operations
        temp_path = temp_file.name

        # Perform your operations
        return load_content(temp_path, extension, loader_extra_config, llm)



def file_to_bytes(filepath):
    """
    Reads a file and returns its content as a bytes object.

    Args:
        filepath (str): The path to the file.

    Returns:
        bytes: The content of the file as a bytes object.
    """
    try:
        with open(filepath, "rb") as f:
            file_content_bytes = f.read()
        return file_content_bytes
    except FileNotFoundError:
        logger.error(f"File not found: {filepath}")
        return None
    except Exception as e:
        logger.error(f"Error reading file {filepath}: {e}")
        return None