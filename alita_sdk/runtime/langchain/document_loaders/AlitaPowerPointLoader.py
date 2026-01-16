import io

from langchain_core.tools import ToolException
from pptx import Presentation
from .utils import perform_llm_prediction_for_image_bytes, create_temp_file
from pptx.enum.shapes import MSO_SHAPE_TYPE
from langchain_core.documents import Document


class AlitaPowerPointLoader:

    def __init__(self, file_path=None, file_content=None, mode=None, **unstructured_kwargs):
        if file_path:
            self.file_path = file_path
        elif file_content:
            self.file_content = file_content
        else:
            raise ToolException("'file_path' or 'file_content' parameter should be provided.")

        self.mode=mode
        self.unstructured_kwargs = unstructured_kwargs
        self.page_number = unstructured_kwargs.get('page_number', None)
        self.extract_images = unstructured_kwargs.get('extract_images', False)
        self.llm = unstructured_kwargs.get('llm', None)
        self.prompt = unstructured_kwargs.get('prompt', "Describe image")
        self.pages_per_chunk = unstructured_kwargs.get('pages_per_chunk', 5)

    def get_content(self):
        if hasattr(self, 'file_path'):
            with open(self.file_path, 'rb') as f:
                prs = Presentation(f)
        elif hasattr(self, 'file_content'):
            prs = Presentation(io.BytesIO(self.file_content))
        pages = []
        if self.page_number is not None:
            pages.append(self.read_pptx_slide(prs.slides[self.page_number - 1], self.page_number))
        else:
            for index, slide in enumerate(prs.slides, start=1):
                pages.append(self.read_pptx_slide(slide, index))
        if self.mode == 'single':
            return "\n".join(pages)
        if self.mode == 'paged':
            return pages
        else:
            raise ToolException(f"Unknown mode value: {self.mode}. Only 'single', 'paged' values allowed.")

    def _extract_table_as_markdown(self, table) -> str:
        """Convert PPTX table to markdown format."""
        if not table.rows:
            return ""
        
        rows = []
        for row in table.rows:
            cells = []
            for cell in row.cells:
                cell_text = cell.text.strip().replace("|", "\\|").replace("\n", " ")
                cells.append(cell_text)
            rows.append("| " + " | ".join(cells) + " |")
        
        if len(rows) > 0:
            # Add header separator after first row
            num_cols = len(table.rows[0].cells)
            header_sep = "| " + " | ".join(["---"] * num_cols) + " |"
            rows.insert(1, header_sep)
        
        return "\n**Table:**\n" + "\n".join(rows) + "\n"

    def _extract_chart_info(self, chart) -> str:
        """Extract data and labels from PPTX chart."""
        result = []
        
        # Extract chart title
        try:
            if chart.has_title and chart.chart_title.has_text_frame:
                title_text = chart.chart_title.text_frame.text.strip()
                if title_text:
                    result.append(f"Chart Title: {title_text}")
        except Exception:
            pass
        
        # Try to extract series data directly from chart.series (works for some chart types)
        try:
            if hasattr(chart, 'series') and chart.series:
                for series in chart.series:
                    series_name = series.name if series.name else "Unnamed Series"
                    values = []
                    categories = []
                    
                    # Try to get values
                    try:
                        if hasattr(series, 'values') and series.values:
                            values = list(series.values)
                    except Exception:
                        pass
                    
                    # Try to get categories from series
                    try:
                        if hasattr(series, 'categories') and series.categories:
                            categories = list(series.categories)
                    except Exception:
                        pass
                    
                    # Build output
                    if categories and values and len(categories) == len(values):
                        data_pairs = [f"{cat}: {val}" for cat, val in zip(categories, values)]
                        result.append(f"Series '{series_name}': {', '.join(data_pairs)}")
                    elif values:
                        result.append(f"Series '{series_name}': {', '.join(str(v) for v in values)}")
                    elif categories:
                        result.append(f"Series '{series_name}' categories: {', '.join(str(c) for c in categories)}")
        except Exception:
            pass
        
        # Fallback: try plots API for bar/line charts
        if not result or (len(result) == 1 and "Chart Title" in result[0]):
            try:
                if hasattr(chart, 'plots') and chart.plots and len(chart.plots) > 0:
                    plot = chart.plots[0]
                    categories = []
                    if hasattr(plot, 'categories') and plot.categories:
                        categories = list(plot.categories)
                        if categories:
                            result.append(f"Categories: {', '.join(str(c) for c in categories)}")
                    
                    # Extract series data from plot
                    for series in plot.series:
                        series_name = series.name if series.name else "Unnamed Series"
                        values = list(series.values) if series.values else []
                        
                        if categories and len(categories) == len(values):
                            data_pairs = [f"{cat}: {val}" for cat, val in zip(categories, values)]
                            result.append(f"Series '{series_name}': {', '.join(data_pairs)}")
                        elif values:
                            result.append(f"Series '{series_name}': {', '.join(str(v) for v in values)}")
            except Exception:
                pass
        
        # Final fallback: parse XML directly for unsupported chart types (e.g., pie3DChart)
        if not result or (len(result) == 1 and "Chart Title" in result[0]):
            try:
                result.extend(self._extract_chart_from_xml(chart))
            except Exception:
                pass
        
        # If we still have no data, add a note
        if not result:
            result.append("(Chart detected - there is no parsed data from this type of chart)")
        
        return "\n**Chart:**\n" + "\n".join(result) + "\n"

    def _extract_chart_from_xml(self, chart) -> list:
        """Extract chart data by parsing the underlying XML directly."""
        result = []
        
        # Get the chart part XML
        chart_part = chart.part
        chart_element = chart_part.element
        
        # Define namespaces used in chart XML
        namespaces = {
            'c': 'http://schemas.openxmlformats.org/drawingml/2006/chart',
            'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
        }
        
        # Find all series (ser) elements
        series_elements = chart_element.findall('.//c:ser', namespaces)
        
        for ser in series_elements:
            series_name = "Unnamed Series"
            categories = []
            values = []
            
            # Extract series name from tx/v or tx/strRef
            tx = ser.find('.//c:tx', namespaces)
            if tx is not None:
                v = tx.find('.//c:v', namespaces)
                if v is not None and v.text:
                    series_name = v.text
            
            # Extract category labels from c:cat
            cat = ser.find('.//c:cat', namespaces)
            if cat is not None:
                # Try strRef first (string references)
                str_cache = cat.find('.//c:strCache', namespaces)
                if str_cache is not None:
                    for pt in str_cache.findall('.//c:pt', namespaces):
                        v = pt.find('c:v', namespaces)
                        if v is not None and v.text:
                            categories.append(v.text)
                
                # Try numRef (numeric references used as categories)
                if not categories:
                    num_cache = cat.find('.//c:numCache', namespaces)
                    if num_cache is not None:
                        for pt in num_cache.findall('.//c:pt', namespaces):
                            v = pt.find('c:v', namespaces)
                            if v is not None and v.text:
                                categories.append(v.text)
            
            # Extract values from c:val
            val = ser.find('.//c:val', namespaces)
            if val is not None:
                num_cache = val.find('.//c:numCache', namespaces)
                if num_cache is not None:
                    for pt in num_cache.findall('.//c:pt', namespaces):
                        v = pt.find('c:v', namespaces)
                        if v is not None and v.text:
                            try:
                                values.append(float(v.text))
                            except ValueError:
                                values.append(v.text)
            
            # Build output
            if categories and values and len(categories) == len(values):
                data_pairs = [f"{cat}: {val}" for cat, val in zip(categories, values)]
                result.append(f"Series '{series_name}': {', '.join(data_pairs)}")
            elif values:
                result.append(f"Series '{series_name}': {', '.join(str(v) for v in values)}")
            elif categories:
                result.append(f"Series '{series_name}' categories: {', '.join(str(c) for c in categories)}")
        
        return result

    def read_pptx_slide(self, slide, index):
        text_content = f'Slide: {index}\n'
        for shape in slide.shapes:
            # Handle tables
            if shape.has_table:
                text_content += self._extract_table_as_markdown(shape.table)
            # Handle charts
            elif shape.has_chart:
                text_content += self._extract_chart_info(shape.chart)
            # Handle images - check multiple ways images can be embedded
            elif self.extract_images and self._is_image_shape(shape):
                try:
                    image_blob = self._get_image_blob(shape)
                    if image_blob:
                        caption = perform_llm_prediction_for_image_bytes(image_blob, self.llm, self.prompt)
                        text_content += "\n**Image Transcript:**\n" + caption + "\n--------------------\n"
                except Exception:
                    pass
            # Handle text frames with hyperlinks
            elif hasattr(shape, "text_frame") and shape.text_frame is not None:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.hyperlink and run.hyperlink.address:
                            link_text = run.text.strip() or "Link"
                            link_url = run.hyperlink.address
                            text_content += f" [{link_text}]({link_url}) "
                        else:
                            text_content += run.text
                text_content += "\n"
        return text_content + "\n"

    def _is_image_shape(self, shape) -> bool:
        """Check if shape contains an image using multiple detection methods."""
        # Method 1: Check shape type
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            return True
        # Method 2: Check if shape has image attribute with blob
        if hasattr(shape, 'image') and shape.image is not None:
            try:
                if shape.image.blob:
                    return True
            except Exception:
                pass
        # Method 3: Check for placeholder with image
        if hasattr(shape, 'placeholder_format') and shape.placeholder_format is not None:
            try:
                if hasattr(shape, 'image') and shape.image is not None:
                    return True
            except Exception:
                pass
        return False

    def _get_image_blob(self, shape) -> bytes:
        """Extract image blob from shape using available methods."""
        # Try direct image access
        if hasattr(shape, 'image') and shape.image is not None:
            try:
                return shape.image.blob
            except Exception:
                pass
        return None

    def load(self):
        content = self.get_content()
        if isinstance(content, str):
            yield Document(page_content=content, metadata={})
        elif isinstance(content, list):
            chunk = []
            chunk_count = 0
            for page_number, page in enumerate(content, start=1):
                chunk.append(page)
                if len(chunk) == self.pages_per_chunk:
                    chunk_content = "\n".join(chunk)
                    yield Document(
                        page_content=chunk_content,
                        metadata={"chunk_number": chunk_count + 1,
                                  "pages_in_chunk": list(range(page_number - len(chunk) + 1, page_number + 1))}
                    )
                    chunk = []
                    chunk_count += 1
            if chunk:
                chunk_content = "\n".join(chunk)
                yield Document(
                    page_content=chunk_content,
                    metadata={"chunk_number": chunk_count + 1,
                              "pages_in_chunk": list(range(len(content) - len(chunk) + 1, len(content) + 1))}
                )